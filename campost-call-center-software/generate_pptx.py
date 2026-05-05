#!/usr/bin/env python3
"""Generate Campost Call Center PowerPoint - Simple language for non-technical audience."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
GOLD = RGBColor(0xC8, 0x96, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
BLUE_LIGHT = RGBColor(0x34, 0x98, 0xDB)


def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return tf


def add_icon_card(slide, left, top, width, height, icon, title, desc, bg_color=WHITE, title_color=DARK_BLUE):
    shape = add_shape_bg(slide, left, top, width, height, bg_color)
    shape.shadow.inherit = False
    # Icon
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(1), Inches(0.6),
                 icon, font_size=28, color=GOLD, bold=True, align=PP_ALIGN.LEFT)
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), Inches(0.5),
                 title, font_size=15, color=title_color, bold=True)
    # Desc
    add_text_box(slide, left + Inches(0.3), top + Inches(1.15), width - Inches(0.6), height - Inches(1.4),
                 desc, font_size=12, color=MED_GRAY)


# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)

# Gold accent bar
add_shape_bg(slide, Inches(0), Inches(3.3), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             'PROPOSITION TECHNIQUE', font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.8),
             'Solution pour le Centre d\'Appels de la CAMPOST', font_size=28, color=GOLD, bold=False, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             'Societe des Postes du Cameroun', font_size=22, color=WHITE, bold=False, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             f'Date : {datetime.date.today().strftime("%d/%m/%Y")}  |  Document Confidentiel',
             font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: SOMMAIRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             'SOMMAIRE', font_size=32, color=WHITE, bold=True)

items_left = [
    '1.  Le Probleme Aujourd\'hui',
    '2.  Qui est la CAMPOST ?',
    '3.  Les Services de la CAMPOST',
    '4.  Ce que les Clients Demandent',
    '5.  Notre Solution',
]
items_right = [
    '6.  Comment ca Marche ?',
    '7.  Les 3 Options Possibles',
    '8.  Calendrier de Realisation',
    '9.  Combien ca Coute ?',
    '10. Prochaines Etapes',
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(4.5), items_left, font_size=20, color=DARK_BLUE, spacing=Pt(14))
add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), items_right, font_size=20, color=DARK_BLUE, spacing=Pt(14))

# ============================================================
# SLIDE 3: LE PROBLEME
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '1. LE PROBLEME AUJOURD\'HUI', font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             'Aujourd\'hui, quand un client de la CAMPOST a un probleme, il n\'a pas de moyen simple pour obtenir de l\'aide.',
             font_size=20, color=DARK_GRAY)

# Problem cards
problems = [
    ('Tel.', 'Un seul numero', 'Les clients appellent le\n+237 2 22 50 70 00 mais\nsouvent personne ne repond\nou la ligne est occupee.'),
    ('Email', 'Reponse lente', 'Les clients ecrivent a\ncs@campostmoney.cm mais\nattendent des jours ou\ndes semaines pour une reponse.'),
    ('Agence', 'Se deplacer', 'Les clients doivent aller\nphysiquement au bureau de\nposte, meme pour une simple\nquestion sur un colis.'),
    ('???', 'Pas de suivi', 'Quand un client se plaint,\nil n\'y a pas de systeme pour\nsuivre sa demande. Elle\npeut se perdre.'),
]
for i, (icon, title, desc) in enumerate(problems):
    left = Inches(0.8 + i * 3.1)
    add_icon_card(slide, left, Inches(2.8), Inches(2.8), Inches(3.2), icon, title, desc, bg_color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
             'Resultat : Des clients mecontents, une image de marque degradee, et des revenus perdus.',
             font_size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: QUI EST LA CAMPOST
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '2. LA CAMPOST EN CHIFFRES', font_size=32, color=WHITE, bold=True)

stats = [
    ('234', 'Bureaux de Poste\nau Cameroun'),
    ('52M$', 'Centre de Donnees\na Yaounde'),
    ('#237#', 'Premier code USSD\nnational en Afrique'),
    ('2-3 Mds', 'FCFA investis dans\nBolamba (e-commerce)'),
]
for i, (num, label) in enumerate(stats):
    left = Inches(0.5 + i * 3.2)
    add_shape_bg(slide, left, Inches(1.8), Inches(2.9), Inches(2.2), DARK_BLUE)
    add_text_box(slide, left, Inches(2.0), Inches(2.9), Inches(0.9),
                 num, font_size=40, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.8), Inches(2.9), Inches(1.0),
                 label, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5),
             'La CAMPOST n\'est plus seulement la Poste. C\'est aussi :\n\n'
             '    CampostMoney  -  Une banque mobile pour epargner, emprunter et transferer de l\'argent\n'
             '    NPSI #237#  -  Un systeme de paiement par telephone (USSD) utilise par les banques\n'
             '    Bolamba (2026)  -  Une plateforme de vente en ligne avec livraison\n'
             '    Camwater  -  Paiement des factures d\'eau via CampostMoney',
             font_size=17, color=DARK_GRAY)

# ============================================================
# SLIDE 5: CE QUE LES CLIENTS DEMANDENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '3. CE QUE LES CLIENTS DEMANDENT LE PLUS', font_size=32, color=WHITE, bold=True)

# Bars representing volume
categories = [
    ('"Ou est mon colis ?"', 'Suivi des colis et courriers', 90, GREEN),
    ('"Mon compte CampostMoney"', 'Solde, transferts, problemes', 75, BLUE_LIGHT),
    ('"Horaires et adresses"', 'Bureaux de poste, tarifs', 65, GOLD),
    ('"J\'ai un probleme #237#"', 'Echecs de paiement USSD', 50, ORANGE),
    ('"Ma facture Camwater"', 'Paiements et confirmations', 40, RGBColor(0x8E, 0x44, 0xAD)),
    ('"Commande Bolamba" (2026)', 'Suivi des commandes en ligne', 30, RGBColor(0x16, 0xA0, 0x85)),
]

for i, (question, desc, pct, color) in enumerate(categories):
    y = Inches(1.6 + i * 0.9)
    add_text_box(slide, Inches(0.8), y, Inches(3.5), Inches(0.4),
                 question, font_size=15, color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.35), Inches(3.5), Inches(0.3),
                 desc, font_size=12, color=MED_GRAY)
    # Bar
    bar_max = Inches(7.5)
    bar_width = int(bar_max * pct / 100)
    add_shape_bg(slide, Inches(4.8), y + Inches(0.08), bar_width, Inches(0.35), color)
    add_text_box(slide, Inches(4.8) + bar_width + Inches(0.15), y, Inches(1), Inches(0.4),
                 f'{pct}%', font_size=14, color=color, bold=True)

# ============================================================
# SLIDE 6: NOTRE SOLUTION - SIMPLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '4. NOTRE SOLUTION : UN CENTRE D\'APPELS MODERNE', font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
             'Un systeme unique ou les agents de la CAMPOST peuvent gerer TOUS les appels et messages des clients.',
             font_size=19, color=DARK_GRAY)

cards = [
    ('Telephone', 'Les clients appellent\nun numero unique.\nLe systeme repond en\nfrancais ou anglais et\ndirige vers le bon agent.', DARK_BLUE),
    ('WhatsApp', 'Les clients ecrivent sur\nWhatsApp. Un robot repond\nautomatiquement aux\nquestions simples (suivi\nde colis, horaires...)', RGBColor(0x25, 0xD3, 0x66)),
    ('SMS', 'Les clients recoivent\ndes notifications par\nSMS : confirmation de\ncolis, rappels, numeros\nde suivi.', BLUE_LIGHT),
    ('Ecran Agent', 'L\'agent voit TOUT sur\nun seul ecran : le colis,\nle compte CampostMoney,\nles paiements #237#.\nPlus besoin de chercher.', GOLD),
]

for i, (title, desc, color) in enumerate(cards):
    left = Inches(0.5 + i * 3.2)
    add_shape_bg(slide, left, Inches(2.5), Inches(2.9), Inches(3.5), color)
    add_text_box(slide, left, Inches(2.7), Inches(2.9), Inches(0.5),
                 title, font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(2.5), Inches(2.5),
                 desc, font_size=14, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
             'Tout est installe sur le centre de donnees de la CAMPOST a Yaounde. Les donnees restent au Cameroun.',
             font_size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: COMMENT CA MARCHE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '5. COMMENT CA MARCHE ? (EN SIMPLE)', font_size=32, color=WHITE, bold=True)

steps = [
    ('1', 'Le client appelle ou\necrit sur WhatsApp', 'Via MTN, Orange\nou CAMTEL'),
    ('2', 'Le systeme vocal\nrepond automatiquement', '"Appuyez 1 pour le\nfrancais, 2 pour\nl\'anglais..."'),
    ('3', 'Questions simples :\nReponse automatique', 'Suivi de colis,\nhoraires, tarifs\n= sans agent'),
    ('4', 'Questions complexes :\nVers un agent', 'L\'agent voit toutes\nles infos du client\nsur son ecran'),
    ('5', 'Le probleme est\nresolu et enregistre', 'Un ticket est cree\npour suivre et ne\nrien oublier'),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.3 + i * 2.6)
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.85), Inches(1.7), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow (except last)
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.2), Inches(1.85), Inches(0.4), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        arrow.line.fill.background()

    add_text_box(slide, left + Inches(0.1), Inches(2.6), Inches(2.3), Inches(1.2),
                 title, font_size=15, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(3.7), Inches(2.3), Inches(1.2),
                 desc, font_size=13, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Bottom benefit
add_shape_bg(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), LIGHT_GRAY)
benefits = [
    'Le client n\'attend plus au telephone - il peut etre rappele',
    'Les questions simples sont traitees automatiquement (60-70% des appels)',
    'L\'agent a toutes les informations sous les yeux, pas besoin de faire attendre le client',
    'Chaque demande est enregistree, rien ne se perd',
]
add_bullet_list(slide, Inches(1), Inches(5.4), Inches(11), Inches(1.6), benefits, font_size=15, color=DARK_BLUE, spacing=Pt(6))

# ============================================================
# SLIDE 8: CANAUX DE COMMUNICATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '6. TOUS LES MOYENS DE CONTACT EN UN SEUL ENDROIT', font_size=28, color=WHITE, bold=True)

channels = [
    ('Appel Telephone', 'Le client appelle,\nle systeme repond et\noriente vers le bon agent.\nFonctionne avec MTN,\nOrange et CAMTEL.', DARK_BLUE),
    ('WhatsApp', 'Le client envoie un\nmessage sur WhatsApp.\nUn robot repond aux\nquestions frequentes.\nL\'agent prend le relais\nsi besoin.', RGBColor(0x25, 0xD3, 0x66)),
    ('SMS', 'Le systeme envoie des\nSMS aux clients :\nconfirmations, numeros\nde suivi, rappels\nde rendez-vous.', BLUE_LIGHT),
    ('USSD #237#', 'Le client tape #237#\nsur son telephone (meme\nsans internet) pour\nverifier un colis ou\nun paiement.', ORANGE),
    ('Email', 'Les demandes par\nemail arrivent dans\nle meme systeme que\nles appels. Rien\nn\'est oublie.', RGBColor(0x8E, 0x44, 0xAD)),
    ('Chat en Ligne', 'Sur le site web\ncampost.cm, les\nvisiteurs peuvent\nchatter en direct\navec un agent.', RGBColor(0x16, 0xA0, 0x85)),
]

for i, (title, desc, color) in enumerate(channels):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.6 + row * 2.9)
    add_shape_bg(slide, left, top, Inches(3.8), Inches(2.6), color)
    add_text_box(slide, left, top + Inches(0.15), Inches(3.8), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.6), Inches(3.2), Inches(1.8),
                 desc, font_size=13, color=WHITE)

# ============================================================
# SLIDE 9: INTEGRATION CAMPOST
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '7. CONNECTE A TOUS LES SYSTEMES CAMPOST', font_size=30, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
             'L\'agent voit tout sur un seul ecran. Il ne perd plus de temps a chercher dans plusieurs systemes.',
             font_size=18, color=DARK_GRAY)

integrations = [
    ('Suivi des Colis (UPU)', 'L\'agent tape le numero du colis\net voit immediatement ou il se\ntrouve : en transit, en douane,\nau bureau de poste, livre.', GREEN),
    ('CampostMoney', 'L\'agent voit le compte du client :\nsolde, derniers transferts, prets\nen cours. Peut aider a resoudre\nles litiges directement.', GOLD),
    ('NPSI #237#', 'Si un paiement USSD a echoue,\nl\'agent voit ce qui s\'est passe\net peut relancer l\'operation\nou rembourser le client.', ORANGE),
    ('Bolamba (2026)', 'Quand la boutique en ligne\nBolamba sera lancee, le centre\nd\'appels sera deja pret pour\ngerer les commandes et retours.', BLUE_LIGHT),
]

for i, (title, desc, color) in enumerate(integrations):
    left = Inches(0.5 + i * 3.2)
    add_shape_bg(slide, left, Inches(2.5), Inches(2.9), Inches(3.8), LIGHT_GRAY)
    # Color bar at top
    add_shape_bg(slide, left, Inches(2.5), Inches(2.9), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=DARK_BLUE, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(2.5), Inches(2.5),
                 desc, font_size=14, color=DARK_GRAY)

# ============================================================
# SLIDE 10: LES 3 OPTIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '8. TROIS OPTIONS POSSIBLES', font_size=32, color=WHITE, bold=True)

options = [
    ('OPTION A\nRECOMMANDEE', GREEN,
     'Solution Sur Mesure\n(Construite pour la CAMPOST)',
     [
         'Installee sur le centre de donnees CAMPOST',
         'Connectee a CampostMoney, #237#, UPU',
         'Fonctionne meme avec internet lent',
         'Les donnees restent au Cameroun',
         'Pas d\'abonnement mensuel eleve',
     ],
     'FCFA 72M (mise en place)\nFCFA 21.5M/an'),
    ('OPTION B', ORANGE,
     'Solution Africaine en Ligne\n(PressOne ou Ameyo)',
     [
         'Plus rapide a mettre en place',
         'Fonctionnalites de base incluses',
         'Mais : abonnement mensuel eleve',
         'Donnees stockees hors du Cameroun',
         'Integration CAMPOST limitee',
     ],
     'FCFA 30-55M (mise en place)\nFCFA 18-45M/an'),
    ('OPTION C', RED,
     'Grande Plateforme Internationale\n(Genesys, NICE)',
     [
         'Marques connues mondialement',
         'Fonctionnalites tres avancees',
         'Mais : extremement cher',
         'Necessite internet tres rapide',
         'Support tres loin (USA/Europe)',
     ],
     'FCFA 90-180M (mise en place)\nFCFA 75-170M/an'),
]

for i, (label, color, subtitle, points, cost) in enumerate(options):
    left = Inches(0.5 + i * 4.2)
    # Header
    add_shape_bg(slide, left, Inches(1.5), Inches(3.8), Inches(0.8), color)
    add_text_box(slide, left, Inches(1.55), Inches(3.8), Inches(0.7),
                 label, font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, left + Inches(0.2), Inches(2.4), Inches(3.4), Inches(0.7),
                 subtitle, font_size=13, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    # Points
    y = Inches(3.2)
    for point in points:
        prefix = '+ ' if not point.startswith('Mais') else '- '
        pcolor = GREEN if not point.startswith('Mais') else RED
        add_text_box(slide, left + Inches(0.2), y, Inches(3.4), Inches(0.35),
                     f'{prefix}{point}', font_size=12, color=pcolor if point.startswith('Mais') else DARK_GRAY)
        y += Inches(0.38)
    # Cost box
    add_shape_bg(slide, left + Inches(0.2), Inches(5.6), Inches(3.4), Inches(0.9), LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.3), Inches(5.65), Inches(3.2), Inches(0.8),
                 cost, font_size=13, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             'Nous recommandons l\'Option A : le meilleur rapport qualite-prix, avec controle total des donnees.',
             font_size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: CALENDRIER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '9. CALENDRIER : 6 MOIS POUR TOUT METTRE EN PLACE', font_size=28, color=WHITE, bold=True)

phases = [
    ('Mois 1', 'PREPARATION', 'Installation des\nequipements sur le\ncentre de donnees\nCAMPOST a Yaounde.', DARK_BLUE),
    ('Mois 2-3', 'CONSTRUCTION', 'Developpement de\nl\'ecran des agents,\ndu systeme vocal\net des rapports.', BLUE_LIGHT),
    ('Mois 3-4', 'CONNEXIONS', 'Branchement avec\nCampostMoney, suivi\ndes colis, #237# et\nWhatsApp.', GOLD),
    ('Mois 5', 'TEST PILOTE', 'Test a Yaounde +\n5 bureaux de poste.\nFormation des agents\nen francais/anglais.', ORANGE),
    ('Mois 6', 'LANCEMENT', 'Deploiement dans\ntous les bureaux.\nLe centre d\'appels\nest operationnel !', GREEN),
]

# Timeline bar
add_shape_bg(slide, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.1), RGBColor(0xDD, 0xDD, 0xDD))

for i, (month, title, desc, color) in enumerate(phases):
    left = Inches(0.5 + i * 2.5)
    # Dot on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.95), Inches(2.85), Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text_box(slide, left, Inches(1.7), Inches(2.3), Inches(0.4),
                 month, font_size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.1), Inches(2.3), Inches(0.4),
                 title, font_size=13, color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(3.6), Inches(2.1), Inches(2.0),
                 desc, font_size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             'Apres le lancement :\n\n'
             '    Mois 7-12 : Ajout de Bolamba (boutique en ligne) quand elle sera prete\n'
             '    En continu : Amelioration du systeme vocal avec intelligence artificielle\n'
             '    Chaque trimestre : Bilan de performance et formation continue des agents',
             font_size=15, color=MED_GRAY)

# ============================================================
# SLIDE 12: COUTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '10. COMBIEN CA COUTE ?', font_size=32, color=WHITE, bold=True)

# Year 1
add_shape_bg(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.5),
             'MISE EN PLACE (Annee 1)', font_size=20, color=DARK_BLUE, bold=True)

items_y1 = [
    ('Equipements et serveurs', '8 000 000 FCFA'),
    ('Systeme telephonique', '5 000 000 FCFA'),
    ('Application pour les agents', '20 000 000 FCFA'),
    ('Connexion CampostMoney', '8 000 000 FCFA'),
    ('Connexion suivi colis', '5 000 000 FCFA'),
    ('Connexion #237# USSD', '6 000 000 FCFA'),
    ('Canal WhatsApp', '4 000 000 FCFA'),
    ('Formation des agents', '4 000 000 FCFA'),
    ('Materiel telephonique', '3 000 000 FCFA'),
    ('Tests et securite', '3 000 000 FCFA'),
    ('Gestion du projet', '6 000 000 FCFA'),
]
y = Inches(2.2)
for label, cost in items_y1:
    add_text_box(slide, Inches(0.8), y, Inches(3.5), Inches(0.3), label, font_size=11, color=DARK_GRAY)
    add_text_box(slide, Inches(4.3), y, Inches(2), Inches(0.3), cost, font_size=11, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    y += Inches(0.3)

add_shape_bg(slide, Inches(0.7), y + Inches(0.05), Inches(5.6), Inches(0.01), GOLD)
add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(3.5), Inches(0.4), 'TOTAL MISE EN PLACE', font_size=14, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(4.0), y + Inches(0.1), Inches(2.3), Inches(0.4), '72 000 000 FCFA', font_size=14, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)

# Annual
add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(3.2), DARK_BLUE)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.5),
             'CHAQUE ANNEE ENSUITE', font_size=20, color=WHITE, bold=True)

items_annual = [
    ('Telephone (appels, SMS)', '6 000 000'),
    ('WhatsApp', '2 000 000'),
    ('Maintenance et support', '8 000 000'),
    ('Hebergement (electricite, reseau)', '2 000 000'),
    ('Licences eventuelles', '1 500 000'),
    ('Formation continue', '2 000 000'),
]
y = Inches(2.3)
for label, cost in items_annual:
    add_text_box(slide, Inches(7.4), y, Inches(3.2), Inches(0.35), label, font_size=13, color=WHITE)
    add_text_box(slide, Inches(10.6), y, Inches(2), Inches(0.35), cost, font_size=13, color=GOLD, align=PP_ALIGN.RIGHT)
    y += Inches(0.4)

add_text_box(slide, Inches(7.4), Inches(4.0), Inches(3), Inches(0.4), 'TOTAL PAR AN', font_size=16, color=WHITE, bold=True)
add_text_box(slide, Inches(10.2), Inches(4.0), Inches(2.4), Inches(0.4), '21 500 000 FCFA', font_size=16, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)

# 3 year TCO
add_shape_bg(slide, Inches(7), Inches(5.0), Inches(5.8), Inches(1.8), GREEN)
add_text_box(slide, Inches(7.2), Inches(5.1), Inches(5.4), Inches(0.4),
             'COUT TOTAL SUR 3 ANS', font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(7.2), Inches(5.6), Inches(5.4), Inches(0.5),
             '143 000 000 FCFA', font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.2), Inches(6.2), Inches(5.4), Inches(0.4),
             'Soit 2 a 3 fois moins cher que les solutions internationales', font_size=13, color=WHITE)

# ============================================================
# SLIDE 13: CONFORMITE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '11. SECURITE ET RESPECT DE LA LOI', font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             'Notre solution respecte la nouvelle loi camerounaise sur la protection des donnees personnelles\n'
             '(Loi N. 2024/017, en vigueur le 23 juin 2026) et les reglements bancaires CEMAC.',
             font_size=17, color=DARK_GRAY)

compliance_cards = [
    ('Donnees au Cameroun', 'Toutes les donnees des clients\nsont stockees au centre de\ndonnees CAMPOST a Yaounde.\nRien ne part a l\'etranger.', GREEN),
    ('Appels Enregistres', 'Les appels concernant\nCampostMoney sont enregistres\ncomme l\'exige le reglement\nbancaire CEMAC (COBAC).', BLUE_LIGHT),
    ('Acces Securise', 'Chaque agent a son propre\ncompte. Les superviseurs\ncontrolent qui accede a\nquelles informations.', DARK_BLUE),
    ('Consentement', 'Le systeme demande l\'accord\ndu client avant d\'enregistrer\nson appel, comme l\'exige\nla loi.', GOLD),
]

for i, (title, desc, color) in enumerate(compliance_cards):
    left = Inches(0.5 + i * 3.2)
    add_shape_bg(slide, left, Inches(2.8), Inches(2.9), Inches(3.0), LIGHT_GRAY)
    add_shape_bg(slide, left, Inches(2.8), Inches(2.9), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=DARK_BLUE, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(3.6), Inches(2.5), Inches(2.0),
                 desc, font_size=14, color=DARK_GRAY)

# ============================================================
# SLIDE 14: POURQUOI NOUS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             '12. POURQUOI NOUS CHOISIR ?', font_size=32, color=WHITE, bold=True)

reasons = [
    ('Nous Connaissons le Cameroun', 'Nous comprenons les reseaux MTN,\nOrange, CAMTEL, le marche bilingue\net l\'environnement FCFA.'),
    ('Nous Connaissons la CAMPOST', 'Nous avons etudie CampostMoney,\nNPSI #237#, UPU et Bolamba en\ndetail pour une integration parfaite.'),
    ('Vos Donnees Restent Chez Vous', 'Tout est installe sur votre propre\ncentre de donnees. Aucune donnee\nclient ne quitte le Cameroun.'),
    ('Prix Juste', 'Notre solution coute 2 a 3 fois\nmoins cher que les alternatives\ninternationales, sans sacrifier\nla qualite.'),
    ('Pret pour l\'Avenir', 'L\'architecture est modulaire :\nBolamba, intelligence artificielle\net nouveaux services s\'ajoutent\nfacilement.'),
    ('Conformite Garantie', 'Loi 2024/017 et reglements\nCOBAC respectes des le\npremier jour de fonctionnement.'),
]

for i, (title, desc) in enumerate(reasons):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.6 + row * 2.8)
    add_shape_bg(slide, left, top, Inches(3.8), Inches(2.5), LIGHT_GRAY)
    add_shape_bg(slide, left, top, Inches(0.08), Inches(2.5), GOLD)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(3.3), Inches(0.5),
                 title, font_size=16, color=DARK_BLUE, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.3), Inches(1.6),
                 desc, font_size=14, color=DARK_GRAY)

# ============================================================
# SLIDE 15: PROCHAINES ETAPES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(3.3), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             'PROCHAINES ETAPES', font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

steps_final = [
    '1.  Reunion de presentation avec l\'equipe CAMPOST',
    '2.  Visite du centre de donnees de Yaounde',
    '3.  Validation des besoins et du cahier des charges',
    '4.  Signature du contrat',
    '5.  Demarrage du projet - le centre d\'appels sera pret en 6 mois',
]
add_bullet_list(slide, Inches(2.5), Inches(2.2), Inches(8), Inches(3.0),
                steps_final, font_size=22, color=WHITE, spacing=Pt(16))

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
             'Ensemble, modernisons le service client de la CAMPOST.',
             font_size=24, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             'Merci  |  Contact : [votre email]  |  [votre telephone]',
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

# Save
output_path = '/Users/user/Documents/hairstyle/hairstyle/campost-call-center-software/Presentation_CAMPOST_Centre_Appels.pptx'
prs.save(output_path)
print(f'PowerPoint saved to: {output_path}')
